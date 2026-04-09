from pptx import Presentation

def generate_ppt(summary):

    prs = Presentation()

    sentences = summary.split(".")

    for s in sentences:

        if len(s.strip()) == 0:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Key Point"
        content.text = s

    prs.save("summary_presentation.pptx")

    return "summary_presentation.pptx"